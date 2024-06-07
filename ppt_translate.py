import argparse
from pptx import Presentation
from deep_translator import GoogleTranslator
from tqdm import tqdm

# Function to translate text in a shape
def translate_shape_text(shape, translator):
    if shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text = run.text
                try:
                    translated_text = translator.translate(text)
                    if translated_text is not None:
                        run.text = translated_text
                    else:
                        print(f"Warning: Translation returned None for text: {text}")
                except Exception as e:
                    print(f"Error translating text: {text}. Error: {e}")

# Function to translate text in a table
def translate_table(table, translator):
    for row in table.rows:
        for cell in row.cells:
            if cell.text_frame:
                translate_shape_text(cell, translator)

# Function to translate text in SmartArt
def translate_smartart(smart_art, translator):
    for shape in smart_art.shapes:
        translate_shape_text(shape, translator)
        if shape.has_smart_art:
            translate_smartart(shape.smart_art, translator)

def main(input_file, source_lang, target_lang):
    # Load the presentation
    prs = Presentation(input_file)

    # Initialize the translator
    translator = GoogleTranslator(source=source_lang, target=target_lang)

    # Calculate the total number of shapes to translate
    total_shapes = sum(len(slide.shapes) for slide in prs.slides)

    # Iterate through each slide in the original presentation with a progress bar
    with tqdm(total=total_shapes, desc="Translating slides", unit="shape") as pbar:
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    translate_shape_text(shape, translator)
                elif shape.has_table:
                    translate_table(shape.table, translator)
                elif hasattr(shape, 'has_smart_art') and shape.has_smart_art:
                    translate_smartart(shape.smart_art, translator)
                pbar.update(1)

    # Save the translated presentation
    output_file = input_file.replace('.pptx', '_translated.pptx')
    prs.save(output_file)
    print(f"Translated presentation saved as {output_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Translate a PowerPoint presentation.')
    parser.add_argument('filename', type=str, help='The path to the PowerPoint file to be translated')
    parser.add_argument('--source', type=str, default='de', help='Source language (default: de)')
    parser.add_argument('--target', type=str, default='en', help='Target language (default: en)')
    args = parser.parse_args()
    main(args.filename, args.source, args.target)
